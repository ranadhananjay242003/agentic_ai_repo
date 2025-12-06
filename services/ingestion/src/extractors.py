"""
Document extractors for various formats
"""
import io
from typing import Tuple, Dict, Any
import logging

logger = logging.getLogger(__name__)

class PDFExtractor:
    def extract(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PDF"""
        try:
            import pdfplumber
            
            pdf_file = io.BytesIO(content)
            text_parts = []
            metadata = {"format": "pdf", "pages": []}
            
            with pdfplumber.open(pdf_file) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text() or ""
                    text_parts.append(page_text)
                    metadata["pages"].append({
                        "page_num": page_num,
                        "char_count": len(page_text)
                    })
            
            full_text = "\n\n".join(text_parts)
            metadata["total_pages"] = len(metadata["pages"])
            
            return full_text, metadata
        
        except ImportError:
            logger.warning("pdfplumber not available, using fallback")
            # Fallback: basic text extraction
            return f"[PDF content placeholder - {len(content)} bytes]", {"format": "pdf", "fallback": True}
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            raise

class DOCXExtractor:
    def extract(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from DOCX"""
        try:
            from docx import Document
            
            doc_file = io.BytesIO(content)
            doc = Document(doc_file)
            
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n\n".join(paragraphs)
            
            metadata = {
                "format": "docx",
                "paragraph_count": len(paragraphs),
                "total_chars": len(full_text)
            }
            
            return full_text, metadata
        
        except ImportError:
            logger.warning("python-docx not available, using fallback")
            return f"[DOCX content placeholder - {len(content)} bytes]", {"format": "docx", "fallback": True}
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            raise

class PPTXExtractor:
    def extract(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PPTX"""
        try:
            from pptx import Presentation
            
            pptx_file = io.BytesIO(content)
            prs = Presentation(pptx_file)
            
            slides_text = []
            metadata = {"format": "pptx", "slides": []}
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text)
                
                slide_text = "\n".join(slide_text_parts)
                slides_text.append(f"[Slide {slide_num}]\n{slide_text}")
                metadata["slides"].append({
                    "slide_num": slide_num,
                    "char_count": len(slide_text)
                })
            
            full_text = "\n\n".join(slides_text)
            metadata["total_slides"] = len(metadata["slides"])
            
            return full_text, metadata
        
        except ImportError:
            logger.warning("python-pptx not available, using fallback")
            return f"[PPTX content placeholder - {len(content)} bytes]", {"format": "pptx", "fallback": True}
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            raise

class CSVExtractor:
    def extract(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from CSV"""
        try:
            import pandas as pd
            
            csv_file = io.BytesIO(content)
            df = pd.read_csv(csv_file)
            
            # Convert to formatted text
            text_parts = [f"CSV Data ({len(df)} rows, {len(df.columns)} columns)\n"]
            text_parts.append(f"Columns: {', '.join(df.columns)}\n")
            text_parts.append(df.to_string())
            
            full_text = "\n".join(text_parts)
            
            metadata = {
                "format": "csv",
                "rows": len(df),
                "columns": list(df.columns),
                "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()}
            }
            
            return full_text, metadata
        
        except ImportError:
            logger.warning("pandas not available, using fallback")
            text = content.decode("utf-8", errors="ignore")
            return text, {"format": "csv", "fallback": True}
        except Exception as e:
            logger.error(f"CSV extraction error: {e}")
            raise
